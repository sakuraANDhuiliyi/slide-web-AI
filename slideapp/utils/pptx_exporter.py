import os
import re
import io
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import markdown
import requests
from urllib.parse import urlparse
from django.conf import settings
from bs4 import BeautifulSoup

def markdown_to_pptx(markdown_content, title="幻灯片"):
    """
    将Markdown格式的幻灯片内容转换为PPTX格式
    
    Args:
        markdown_content (str): Markdown格式的幻灯片内容
        title (str): 幻灯片标题
    
    Returns:
        BytesIO: PPTX文件的二进制内容
    """
    prs = Presentation()
    
    # 设置幻灯片大小为16:9比例
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # 解析幻灯片结构（水平、垂直和渐变）
    slides_data = parse_markdown_slides(markdown_content)
    
    # 处理每个幻灯片
    for slide_data in slides_data:
        # 选择布局（标题幻灯片或标题和内容）
        if slide_data['is_title_slide']:
            slide_layout = prs.slide_layouts[0]  # 标题幻灯片
        else:
            slide_layout = prs.slide_layouts[1]  # 标题和内容幻灯片
        
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        if slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = slide_data['title']
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(40)
                    run.font.bold = True
                    run.font.color.rgb = RGBColor(0, 0, 0)
        
        # 设置内容
        if slide_data['content'] and hasattr(slide.shapes, 'placeholders'):
            try:
                # 尝试使用内容占位符
                content_placeholder = None
                for shape in slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.type == 1:  # 1表示内容占位符
                        content_placeholder = shape
                        break
                
                if not content_placeholder and len(slide.shapes.placeholders) > 1:
                    content_placeholder = slide.shapes.placeholders[1]
                
                if content_placeholder:
                    # 将Markdown转换为HTML
                    html_content = markdown.markdown(slide_data['content'], extensions=['tables', 'fenced_code'])
                    
                    # 解析HTML获取纯文本
                    soup = BeautifulSoup(html_content, 'html.parser')
                    text = soup.get_text('\n')
                    
                    # 添加内容
                    tf = content_placeholder.text_frame
                    tf.clear()  # 清除默认文本
                    
                    # 添加段落
                    for paragraph_text in text.split('\n'):
                        if paragraph_text.strip():
                            p = tf.add_paragraph()
                            p.text = paragraph_text
                            p.font.size = Pt(24)
                
                # 处理图片
                if slide_data['images']:
                    for img_data in slide_data['images']:
                        try:
                            # 下载图片
                            image_path = download_image(img_data['url'])
                            if image_path:
                                # 添加图片到幻灯片
                                left = Inches(1)
                                top = Inches(3)  # 在内容之后显示
                                width = Inches(8)  # 宽度自适应
                                slide.shapes.add_picture(image_path, left, top, width=width)
                                # 删除临时文件
                                os.remove(image_path)
                        except Exception as e:
                            print(f"添加图片时出错: {str(e)}")
            except Exception as e:
                print(f"添加内容时出错: {str(e)}")
    
    # 保存为二进制内容
    pptx_data = io.BytesIO()
    prs.save(pptx_data)
    pptx_data.seek(0)
    
    return pptx_data

def parse_markdown_slides(markdown_content):
    """
    解析Markdown内容，提取幻灯片结构
    """
    slides_data = []
    
    # 分割水平幻灯片
    h_slides = re.split(r'\n\s*---\s*\n', markdown_content)
    
    for h_slide in h_slides:
        # 检查是否为垂直幻灯片组
        if '\n----\n' in h_slide:
            # 分割垂直幻灯片
            v_slides = re.split(r'\n\s*----\s*\n', h_slide)
            
            for v_idx, v_slide in enumerate(v_slides):
                # 检查是否为渐变幻灯片组
                if '\n++++\n' in v_slide:
                    # 分割渐变幻灯片
                    a_slides = re.split(r'\n\s*\+\+\+\+\s*\n', v_slide)
                    
                    for a_idx, a_slide in enumerate(a_slides):
                        slide_data = process_slide_content(a_slide)
                        slides_data.append(slide_data)
                else:
                    slide_data = process_slide_content(v_slide)
                    slides_data.append(slide_data)
        else:
            # 普通水平幻灯片
            slide_data = process_slide_content(h_slide)
            slides_data.append(slide_data)
    
    return slides_data

def process_slide_content(slide_content):
    """
    处理单个幻灯片内容，提取标题、内容和图片
    """
    slide_data = {
        'title': '',
        'content': '',
        'images': [],
        'is_title_slide': False
    }
    
    # 提取标题
    title_match = re.search(r'^\s*(##+)\s+(.+)$', slide_content, re.MULTILINE)
    if title_match:
        heading_level = len(title_match.group(1))
        slide_data['title'] = title_match.group(2).strip()
        slide_data['is_title_slide'] = (heading_level == 1)
        
        # 移除标题行
        slide_content = re.sub(r'^\s*##+\s+.+$', '', slide_content, 1, re.MULTILINE)
    
    # 提取图片
    img_matches = re.finditer(r'!\[(.*?)\]\((.*?)\)', slide_content)
    for match in img_matches:
        alt_text = match.group(1)
        img_url = match.group(2)
        slide_data['images'].append({
            'alt': alt_text,
            'url': img_url
        })
        
        # 移除图片标记
        slide_content = slide_content.replace(match.group(0), '')
    
    # 剩余内容作为幻灯片正文
    slide_data['content'] = slide_content.strip()
    
    return slide_data

def download_image(image_url):
    """
    下载图片并返回本地路径
    """
    try:
        # 如果是相对路径，转换为绝对路径
        if image_url.startswith('/'):
            # 去掉开头的斜杠以获取相对路径
            relative_path = image_url[1:] if image_url.startswith('/') else image_url
            image_path = os.path.join(settings.BASE_DIR, relative_path)
            if os.path.exists(image_path):
                return image_path
        
        # 下载图片
        response = requests.get(image_url, stream=True, verify=False, timeout=10)
        if response.status_code == 200:
            # 创建临时文件
            fd, temp_path = tempfile.mkstemp(suffix='.jpg')
            os.close(fd)
            
            # 保存图片到临时文件
            with open(temp_path, 'wb') as f:
                for chunk in response.iter_content(1024):
                    f.write(chunk)
            
            return temp_path
    except Exception as e:
        print(f"下载图片失败: {str(e)}")
    
    return None 